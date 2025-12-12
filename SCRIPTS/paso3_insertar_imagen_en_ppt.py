from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

# Rutas
carpeta_imgs = "/Users/Arturo/AGRICULTURA/INFORMES/Img_Temporales"
ruta_pptx = "/Users/Arturo/AGRICULTURA/INFORMES/Avances Fertilizantes 2025.pptx"

# Lista completa de imágenes y sus slides (índice = slide - 1)
imagenes = [
    ("avance_nacional.png", 1),
    ("guerrero_avances.png", 9),
    ("durango_avances.png", 7),
    ("michoacan_avances.png", 11),
    ("morelos_avances.png", 13),
    ("tlaxcala_avances.png", 15),
    ("veracruz_avances.png", 18),
    ("chiapas_avances.png", 20),
    ("tabasco_avances.png", 23),
    ("colima_avances.png", 24),
    ("edomex_avances.png", 25),
    ("cdmx_avanes.png", 26),
    ("guanajuato_avances.png", 27),
    ("jalisco_avances.png", 28),
    ("puebla_avances.png", 29),
    ("hidalgo_avances.png", 30),
    ("campeche_avances.png", 31),
    ("yucatan_avances.png", 32),
    ("quintana_roo_avances.png", 33),
    ("oaxaca_avances.png", 34)
]

# Validar presentación
if not os.path.exists(ruta_pptx):
    print("❌ No se encontró la presentación.")
else:
    ppt = Presentation(ruta_pptx)

    for nombre_archivo, slide_idx in imagenes:
        ruta_img = os.path.join(carpeta_imgs, nombre_archivo)

        if not os.path.exists(ruta_img):
            print(f"❌ No se encontró la imagen: {ruta_img}")
            continue

        if len(ppt.slides) <= slide_idx:
            print(f"❌ La presentación no tiene un slide en el índice {slide_idx + 1}.")
            continue

        slide = ppt.slides[slide_idx]

        # 🧽 Eliminar imágenes anteriores del slide
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sp = shape._element
                sp.getparent().remove(sp)

        # 📌 Insertar imagen con posición exacta
        slide.shapes.add_picture(
            ruta_img,
            Cm(1.15),
            Cm(4.03),
            width=Cm(31.57),
            height=Cm(11)
        )

        print(f"✅ Imagen '{nombre_archivo}' insertada en slide {slide_idx + 1}.")

    ppt.save(ruta_pptx)
    print("💾 Presentación actualizada correctamente.")
