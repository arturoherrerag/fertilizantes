import xlwings as xw
from pptx import Presentation
from pptx.util import Inches
import os
import shutil

# 📁 Archivos
ruta_excel = "/Users/Arturo/AGRICULTURA/fertilizantes_principal_2025.xlsx"
hoja = "Avances Generales"
rango = "B4:Q17"
imagen_temp = "/Users/Arturo/AGRICULTURA/imagen_temp_avance.png"
ruta_presentacion = "/Users/Arturo/AGRICULTURA/INFORMES/Avances Fertilizantes 2025_28032025.pptx"

# 🧩 Abrir Excel y copiar el rango como imagen
app = xw.App(visible=False)
wb = xw.Book(ruta_excel)
sht = wb.sheets[hoja]
sht.range(rango).api.CopyPicture(Appearance=1, Format=2)

# 🖼️ Pegar en hoja temporal
wb_temp = xw.Book()
ws_temp = wb_temp.sheets[1]
ws_temp.api.Paste()
shape = ws_temp.api.Shapes(1)
shape.Export(imagen_temp, 2)
wb_temp.close()
wb.close()
app.quit()
print(f"✅ Imagen creada: {imagen_temp}")

# 🖼️ Abrir presentación y duplicar slide base
ppt = Presentation(ruta_presentacion)
slide_base = ppt.slides[0]  # toma el primero como plantilla
slide_layout = slide_base.slide_layout
slide_nuevo = ppt.slides.add_slide(slide_layout)

# 🔁 Copiar shapes del slide base
for shape in slide_base.shapes:
    el = shape.element
    new_el = el.clone()
    slide_nuevo.shapes._spTree.insert_element_before(new_el, 'p:extLst')

# 🖼️ Agregar la imagen en posición fija
slide_nuevo.shapes.add_picture(imagen_temp, Inches(1), Inches(1.6), height=Inches(4.8))

# 💾 Guardar presentación
ppt.save(ruta_presentacion)
print(f"✅ Slide duplicado con imagen en: {ruta_presentacion}")
