import xlwings as xw
from PIL import ImageGrab
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import time

# --- CONFIGURACIONES ---
ruta_excel = "/Users/Arturo/AGRICULTURA/INFORMES/INFORME_2025.xlsm"
ruta_pptx = "/Users/Arturo/AGRICULTURA/INFORMES/Avances Fertilizantes 2025.pptx"
carpeta_imgs = "/Users/Arturo/AGRICULTURA/INFORMES/Img_Temporales"
os.makedirs(carpeta_imgs, exist_ok=True)

# --- LISTA DE IMÁGENES DESDE EXCEL ---
imagenes_excel = [
    ("Avances Generales", "avance_nacional.png", "B4:Q17", 1, 1.15, 3.69, 31.57),
    ("En entregas", "en_entregas.png", "B3:K20", 5, 0.98, 3.34, 31.93),
    ("En entregas", "en_entregas_2.png", "B22:K39", 6, 0.97, 3.34, 31.93),
]

# --- LISTA DE IMÁGENES YA GENERADAS CON PYTHON ---
imagenes_python = [
    ("grafico_envios_diarios_powerbi.png", 2, 0.43, 3.15, 33),
    ("grafico_entregas_diarias.png", 3, 2.43, 2.83, 29),
    ("grafico_comparativo_abasto_entregas.png", 4, 2.93, 3.34, 28),
]

# --- PARTE 1: EXPORTAR IMÁGENES DESDE EXCEL ---
print("📤 Exportando imágenes desde Excel...")

app = xw.App(visible=False)
wb = xw.Book(ruta_excel)

def exportar_imagen(nombre_hoja, nombre_archivo, rango):
    try:
        macro = wb.macro("ExportarImagenDeHojaRango")
        macro(nombre_hoja, rango)

        time.sleep(1.5)
        imagen = ImageGrab.grabclipboard()
        if imagen:
            ruta_img = os.path.join(carpeta_imgs, nombre_archivo)
            imagen.save(ruta_img, "PNG")
            print(f"✅ Imagen de '{nombre_hoja}' guardada como: {ruta_img}")
        else:
            print(f"❌ No se encontró imagen en portapapeles para '{nombre_hoja}'.")
    except Exception as e:
        print(f"❌ Error con la hoja '{nombre_hoja}': {e}")

for hoja, archivo, rango, _, _, _, _ in imagenes_excel:
    exportar_imagen(hoja, archivo, rango)

wb.close()
app.quit()

# --- PARTE 2: INSERTAR IMÁGENES EN POWERPOINT ---
print("\n📥 Insertando imágenes en PowerPoint...")

if not os.path.exists(ruta_pptx):
    print("❌ No se encontró la presentación.")
else:
    ppt = Presentation(ruta_pptx)

    # Combinar ambas listas de imágenes
    imagenes_todas = [
        *[(archivo, slide, izquierda, arriba, ancho) for _, archivo, _, slide, izquierda, arriba, ancho in imagenes_excel],
        *imagenes_python
    ]

    for archivo, slide_idx, izquierda_cm, arriba_cm, ancho_cm in imagenes_todas:
        ruta_img = os.path.join(carpeta_imgs, archivo)

        if not os.path.exists(ruta_img):
            print(f"❌ No se encontró la imagen: {ruta_img}")
            continue

        if len(ppt.slides) <= slide_idx:
            print(f"❌ La presentación no tiene un slide en el índice {slide_idx + 1}.")
            continue

        slide = ppt.slides[slide_idx]

        # Eliminar imágenes anteriores del slide
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sp = shape._element
                sp.getparent().remove(sp)

        # Insertar nueva imagen
        slide.shapes.add_picture(
            ruta_img,
            Cm(izquierda_cm),
            Cm(arriba_cm),
            width=Cm(ancho_cm)
        )

        print(f"✅ Imagen '{archivo}' insertada en slide {slide_idx + 1}.")

    ppt.save(ruta_pptx)
    print("\n💾 Presentación actualizada correctamente.")