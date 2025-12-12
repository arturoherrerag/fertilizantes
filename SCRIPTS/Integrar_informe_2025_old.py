import xlwings as xw
from PIL import ImageGrab
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import time

# --- CONFIGURACIONES ---

ruta_excel = "/Users/Arturo/AGRICULTURA/INFORMES/INFORME_2025.xlsm"
ruta_pptx = "/Users/Arturo/AGRICULTURA/INFORMES/Avances Fertilizantes 2025.pptx"
carpeta_imgs = "/Users/Arturo/AGRICULTURA/INFORMES/Img_Temporales"
os.makedirs(carpeta_imgs, exist_ok=True)

# Lista de hojas, archivos, rangos y slides
hojas_con_info = [
    ("Avances Generales", "avance_nacional.png", "B4:Q17", 1),
    ("Guerrero Avances", "guerrero_avances.png", "B5:Q18", 9),
    ("Durango Avances", "durango_avances.png", "B5:Q18", 7),
    ("Michoacán Avances", "michoacan_avances.png", "B5:Q18", 11),
    ("Morelos Avances", "morelos_avances.png", "B5:Q18", 13),
    ("Tlaxcala Avances", "tlaxcala_avances.png", "B5:Q18", 15),
    ("Veracruz Avances", "veracruz_avances.png", "B5:Q18", 18),
    ("Chiapas Avances", "chiapas_avances.png", "B5:Q18", 20),
    ("Tabasco Avances", "tabasco_avances.png", "B5:Q18", 23),
    ("Colima Avances", "colima_avances.png", "B5:Q18", 24),
    ("Edomex Avances", "edomex_avances.png", "B5:Q18", 25),
    ("Cdmx Avances", "cdmx_avanes.png", "B5:Q18", 26),
    ("Guanajuato Avances", "guanajuato_avances.png", "B5:Q18", 27),
    ("Jalisco Avances", "jalisco_avances.png", "B5:Q18", 28),
    ("Puebla Avances", "puebla_avances.png", "B5:Q18", 29),
    ("Hidalgo Avances", "hidalgo_avances.png", "B5:Q18", 30),
    ("Campeche Avances", "campeche_avances.png", "B5:Q18", 31),
    ("Yucatán Avances", "yucatan_avances.png", "B5:Q18", 32),
    ("Quintana Roo Avances", "quintana_roo_avances.png", "B5:Q18", 33),
    ("Oaxaca Avances", "oaxaca_avances.png", "B5:Q18", 34)
]

# --- PARTE 1: Exportar imágenes desde Excel ---

print("📤 Exportando imágenes desde Excel...")

app = xw.App(visible=False)
wb = xw.Book(ruta_excel)

def exportar_imagen(nombre_hoja, nombre_archivo, rango):
    try:
        macro = wb.macro("ExportarImagenDeHojaRango")
        macro(nombre_hoja, rango)

        time.sleep(1.5)
        imagen = ImageGrab.grabclipboard()
        if imagen:
            ruta_img = os.path.join(carpeta_imgs, nombre_archivo)
            imagen.save(ruta_img, "PNG")
            print(f"✅ Imagen de '{nombre_hoja}' guardada como: {ruta_img}")
        else:
            print(f"❌ No se encontró imagen en portapapeles para '{nombre_hoja}'.")

    except Exception as e:
        print(f"❌ Error con la hoja '{nombre_hoja}': {e}")

for hoja, archivo, rango, _ in hojas_con_info:
    exportar_imagen(hoja, archivo, rango)

wb.close()
app.quit()

# --- PARTE 2: Insertar imágenes en PowerPoint ---

print("\n📥 Insertando imágenes en PowerPoint...")

if not os.path.exists(ruta_pptx):
    print("❌ No se encontró la presentación.")
else:
    ppt = Presentation(ruta_pptx)

    for _, archivo, _, slide_idx in hojas_con_info:
        ruta_img = os.path.join(carpeta_imgs, archivo)

        if not os.path.exists(ruta_img):
            print(f"❌ No se encontró la imagen: {ruta_img}")
            continue

        if len(ppt.slides) <= slide_idx:
            print(f"❌ La presentación no tiene un slide en el índice {slide_idx + 1}.")
            continue

        slide = ppt.slides[slide_idx]

        # Eliminar imágenes anteriores del slide
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sp = shape._element
                sp.getparent().remove(sp)

        # Insertar nueva imagen
        slide.shapes.add_picture(
            ruta_img,
            Cm(1.15),
            Cm(4.03),
            width=Cm(31.57),
            height=Cm(11)
        )

        print(f"✅ Imagen '{archivo}' insertada en slide {slide_idx + 1}.")

    ppt.save(ruta_pptx)
    print("\n💾 Presentación actualizada correctamente.")
